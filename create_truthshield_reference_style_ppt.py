from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUT = BASE_DIR / "TruthShield_Tomorrow_Presentation.pptx"

SLIDES = [
    {
        "title": "Introduction",
        "subtitle": "TruthShield: AI Deepfake Fraud Detection System",
        "bullets": [
            "TruthShield is a web-based deepfake fraud detection project.",
            "Deepfake media can create fake images, videos, and audio that appear real.",
            "It can be misused for misinformation, impersonation, and digital fraud.",
            "TruthShield helps users verify suspicious media through a web-based system.",
        ],
        "kind": "intro",
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Manual media verification is difficult for normal users.",
            "Fake media spreads quickly through social media and messaging platforms.",
            "Realistic AI-generated content can create trust, privacy, and security issues.",
            "A simple detection platform is needed to support safer digital communication.",
        ],
    },
    {
        "title": "Objective",
        "bullets": [
            "To develop a user-friendly deepfake fraud detection system.",
            "To allow users to upload image, video, or audio files for analysis.",
            "To classify uploaded media as Authenticated or Manipulated.",
            "To maintain user detection history and admin-side monitoring.",
        ],
    },
    {
        "title": "Existing System",
        "bullets": [
            "Many existing systems depend on manual checking.",
            "Some tools support only one type of media.",
            "Advanced systems may require high hardware resources.",
            "Poor-quality media can reduce detection reliability.",
        ],
    },
    {
        "title": "Proposed System",
        "bullets": [
            "TruthShield provides a complete web-based media verification workflow.",
            "Users can register, log in, upload media, and view detection results.",
            "The backend analyses the uploaded file and returns a clear output.",
            "SQLite stores user details and detection history records.",
        ],
    },
    {
        "title": "Technologies Used",
        "bullets": [
            "Frontend: HTML, CSS, JavaScript",
            "Backend: Python Flask",
            "Database: SQLite",
            "Media processing: OpenCV and NumPy",
            "Development tool: Visual Studio Code",
        ],
    },
    {
        "title": "System Architecture",
        "bullets": [
            "User Interface",
            "Flask Backend",
            "Media Analysis Module",
            "SQLite Database",
            "Detection Result",
        ],
        "kind": "architecture",
    },
    {
        "title": "Working Process",
        "bullets": [
            "User registers or logs in to the system.",
            "User uploads suspicious image, video, or audio media.",
            "Backend identifies the media type and preprocesses the file.",
            "Detection logic analyses the file and generates the result.",
            "Result is displayed and saved in detection history.",
        ],
    },
    {
        "title": "Dataset Description",
        "bullets": [
            "Dataset contains real and manipulated media samples.",
            "Videos are arranged into Authenticated and Manipulated categories.",
            "Frames can be extracted from videos for model training and testing.",
            "A larger balanced dataset can improve future accuracy.",
        ],
    },
    {
        "title": "Deepfake Detection Techniques",
        "bullets": [
            "Face visibility and facial consistency checking",
            "Lighting and scene stability analysis",
            "Texture and frame irregularity screening",
            "Motion continuity checking for video files",
            "Confidence score and signal-based explanation",
        ],
    },
    {
        "title": "Model Implementation",
        "bullets": [
            "Uploaded media is received by the Python Flask backend.",
            "Image and video files are preprocessed before analysis.",
            "Video files are checked using selected sampled frames.",
            "Model output is combined with visual screening signals.",
            "Final output is generated as Authenticated or Manipulated.",
        ],
    },
    {
        "title": "Features",
        "bullets": [
            "Professional home page and navigation",
            "User registration and login",
            "Media upload and detection page",
            "Confidence and signal breakdown output",
            "Detection history and admin dashboard",
        ],
    },
    {
        "title": "Advantages",
        "bullets": [
            "Simple and easy-to-use interface",
            "Supports image, video, and audio upload workflow",
            "Helps create awareness about deepfake fraud",
            "Stores detection history for future reference",
            "Provides a strong base for future AI model improvement",
        ],
    },
    {
        "title": "Limitations",
        "bullets": [
            "Detection result depends on media quality and clarity.",
            "Low lighting, blur, and compression may affect the output.",
            "Accuracy can be improved with larger and more balanced datasets.",
            "The current system is a project-level detector, not a forensic replacement.",
        ],
    },
    {
        "title": "Future Enhancements",
        "bullets": [
            "Train the model using larger deepfake datasets.",
            "Improve image, video, and audio detection accuracy.",
            "Add real-time webcam-based detection.",
            "Deploy the system as a full web or mobile application.",
            "Add advanced admin analytics and explainable AI reports.",
        ],
    },
    {
        "title": "Conclusion",
        "bullets": [
            "TruthShield provides a practical solution for deepfake fraud awareness.",
            "The project combines frontend, backend, database, and detection workflow.",
            "It helps users analyse suspicious media in a simple and clear way.",
            "The system can be improved further with stronger models and larger datasets.",
        ],
    },
    {
        "title": "References",
        "bullets": [
            "FaceForensics++: Learning to Detect Manipulated Facial Images",
            "Deepfake Detection Challenge Dataset",
            "Deepfake Video Detection Using Recurrent Neural Networks",
            "Python Flask, SQLite, OpenCV, HTML, CSS, and JavaScript documentation",
        ],
    },
    {
        "title": "Thank You",
        "subtitle": "Questions and Answers",
        "bullets": ["TruthShield: AI Deepfake Fraud Detection System"],
        "kind": "thanks",
    },
]


def add_textbox(slide, text, left, top, width, height, size=24, bold=False, color=(22, 33, 43), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_bullets(slide, bullets, left, top, width, height, size=22):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(35, 35, 35)
        p.space_after = Pt(10)
        p.margin_left = Inches(0.25)
    return box


def add_background(slide, number):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 249, 244)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(14, 47, 55)
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.57), Inches(0.45), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(11, 114, 133)
    accent.line.fill.background()
    add_textbox(slide, f"{number:02}", Inches(11.85), Inches(0.43), Inches(0.6), Inches(0.25), 12, True, (95, 109, 122), PP_ALIGN.RIGHT)
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(7.08), Inches(11.6), Inches(0.01))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(207, 216, 220)
    footer.line.fill.background()
    add_textbox(slide, "AI Deepfake Fraud Detection | Project Presentation", Inches(0.65), Inches(7.18), Inches(5.8), Inches(0.2), 9, False, (95, 109, 122))


def add_title_slide(slide, data):
    add_textbox(slide, "TRUTHSHIELD", Inches(0.75), Inches(1.0), Inches(3.5), Inches(0.4), 14, True, (11, 114, 133))
    add_textbox(slide, data["title"], Inches(0.75), Inches(1.65), Inches(7.6), Inches(0.9), 46, True, (23, 33, 43))
    add_textbox(slide, data["subtitle"], Inches(0.78), Inches(2.65), Inches(7.1), Inches(0.45), 22, False, (95, 109, 122))
    add_bullets(slide, data["bullets"], Inches(1.05), Inches(3.6), Inches(7.0), Inches(1.3), 22)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.45), Inches(2.7), Inches(2.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(16, 42, 51)
    shape.line.fill.background()
    add_textbox(slide, "AI\nMEDIA\nVERIFY", Inches(9.35), Inches(2.02), Inches(2.4), Inches(1.4), 34, True, (255, 255, 255), PP_ALIGN.CENTER)


def add_intro_slide(slide, data):
    add_textbox(slide, "TRUTHSHIELD", Inches(0.75), Inches(0.82), Inches(3.5), Inches(0.35), 14, True, (11, 114, 133))
    add_textbox(slide, data["title"], Inches(0.75), Inches(1.28), Inches(7.6), Inches(0.65), 38, True, (23, 33, 43))
    add_textbox(slide, data["subtitle"], Inches(0.78), Inches(2.0), Inches(7.1), Inches(0.35), 18, False, (95, 109, 122))
    add_bullets(slide, data["bullets"], Inches(1.05), Inches(2.85), Inches(7.8), Inches(2.7), 21)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.35), Inches(2.0), Inches(2.35), Inches(2.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(16, 42, 51)
    shape.line.fill.background()
    add_textbox(slide, "AI\nMEDIA\nVERIFY", Inches(9.52), Inches(2.48), Inches(2.0), Inches(1.25), 28, True, (255, 255, 255), PP_ALIGN.CENTER)


def add_normal_slide(slide, data):
    add_textbox(slide, data["title"], Inches(0.75), Inches(0.95), Inches(8.8), Inches(0.75), 34, True, (23, 33, 43))
    add_bullets(slide, data["bullets"], Inches(1.0), Inches(2.1), Inches(9.7), Inches(4.1), 22)


def add_architecture_slide(slide, data):
    add_textbox(slide, data["title"], Inches(0.75), Inches(0.95), Inches(8.8), Inches(0.75), 34, True, (23, 33, 43))
    add_bullets(slide, data["bullets"], Inches(0.95), Inches(1.95), Inches(3.8), Inches(3.4), 20)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(1.75), Inches(6.8), Inches(4.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(11, 114, 133)
    box.line.width = Pt(1.5)
    add_textbox(slide, "Space for System Architecture Diagram", Inches(5.55), Inches(3.65), Inches(5.8), Inches(0.5), 20, True, (95, 109, 122), PP_ALIGN.CENTER)


def add_thanks_slide(slide, data):
    add_textbox(slide, "Thank You", Inches(0.9), Inches(2.0), Inches(7.0), Inches(0.9), 54, True, (23, 33, 43))
    add_textbox(slide, data["subtitle"], Inches(0.95), Inches(3.05), Inches(5.2), Inches(0.45), 26, True, (11, 114, 133))
    add_textbox(slide, data["bullets"][0], Inches(0.95), Inches(3.8), Inches(6.4), Inches(0.35), 20, False, (95, 109, 122))
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(1.9), Inches(2.7), Inches(2.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(233, 243, 245)
    box.line.color.rgb = RGBColor(183, 221, 228)
    add_textbox(slide, "Q&A", Inches(9.25), Inches(2.72), Inches(1.8), Inches(0.65), 46, True, (11, 114, 133), PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for idx, data in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        add_background(slide, idx)
        if data.get("kind") == "cover":
            add_title_slide(slide, data)
        elif data.get("kind") == "intro":
            add_intro_slide(slide, data)
        elif data.get("kind") == "architecture":
            add_architecture_slide(slide, data)
        elif data.get("kind") == "thanks":
            add_thanks_slide(slide, data)
        else:
            add_normal_slide(slide, data)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
